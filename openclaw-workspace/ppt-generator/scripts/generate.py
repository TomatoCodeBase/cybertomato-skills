# -*- coding: utf-8 -*-
"""
PPT v2 - 小七姐风格复刻
黑/白双色调 + 橙色强调 + 极简扁平 + 大留白
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Color Palette (小七姐风格) ===
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xf8, 0xf8, 0xf8)
ORANGE = RGBColor(0xe8, 0x68, 0x00)
ORANGE_LIGHT = RGBColor(0xf0, 0x68, 0x00)
ORANGE_BG = RGBColor(0xfd, 0xf0, 0xe2)  # Very light orange for highlights
DARK_TEXT = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xc0, 0xc0, 0xc0)
CARD_DARK = RGBColor(0x18, 0x18, 0x18)
CARD_DARK2 = RGBColor(0x22, 0x22, 0x22)
BLUE = RGBColor(0x38, 0x50, 0xd8)

FONT = 'Microsoft YaHei'

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return txBox

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, color=ORANGE, thickness=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_multi(slide, left, top, width, height, lines, size=16, spacing=1.4):
    """lines = [(text, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, clr, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = FONT
        p.space_after = Pt(size * 0.4)
    return txBox

def add_page_num(slide, num, total, dark_bg=False):
    color = WHITE if dark_bg else LIGHT_GRAY
    add_text(slide, Inches(12.5), Inches(7), Inches(0.8), Inches(0.4),
        f'{num}', size=10, color=color, align=PP_ALIGN.RIGHT)

def add_footer_brand(slide, dark_bg=False):
    """小七姐风格：底部品牌标识"""
    color = RGBColor(0x50, 0x50, 0x50) if not dark_bg else RGBColor(0x40, 0x40, 0x40)
    add_text(slide, Inches(0.5), Inches(7), Inches(3), Inches(0.4),
        '赛博番茄  @2026', size=10, color=color)


# ===================================================================
# SLIDE 1: COVER (黑底)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)

# 顶部橙色装饰线
add_line(slide, Inches(0.5), Inches(0.5), Inches(1.5))

# 副标题（小字）
add_text(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5),
    '番茄 · 赛博番茄 · 2026.03', size=14, color=GRAY)

# 主标题（超大）
add_text(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(2),
    '一个人就是\n一家内容公司', size=52, color=WHITE, bold=True)

# 橙色装饰线
add_line(slide, Inches(0.5), Inches(5), Inches(3), ORANGE)

# 副标题
add_text(slide, Inches(0.5), Inches(5.3), Inches(10), Inches(0.8),
    '从零代码到全自动化的公众号实战', size=22, color=ORANGE)


# ===================================================================
# SLIDE 2: COZE - 零门槛 (白底)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.3),
    '01', size=12, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(0.8), Inches(10), Inches(0.8),
    '零门槛：扣子编程', size=36, color=DARK_TEXT, bold=True)

add_line(slide, Inches(0.5), Inches(1.8), Inches(1.5))

add_text(slide, Inches(0.5), Inches(2.2), Inches(10), Inches(0.6),
    '不用写一行代码，用扣子搭建AI自动化公众号', size=18, color=GRAY)

# Flow
flow_steps = ['选题', '标题', '正文', '配图', '摘要']
for i, step in enumerate(flow_steps):
    x = Inches(0.8 + i * 2.4)
    add_rounded(slide, x, Inches(3.2), Inches(2), Inches(0.7), ORANGE if i == 0 else RGBColor(0xf0,0xf0,0xf0))
    add_text(slide, x, Inches(3.25), Inches(2), Inches(0.6),
        step, size=18, color=WHITE if i == 0 else DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        add_text(slide, Inches(0.8 + i * 2.4 + 2), Inches(3.15), Inches(0.4), Inches(0.7),
            '→', size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 3 cards
cards = [
    ('不写一行代码', '扣子长期计划 + 技能商店\n拖拽式搭建，人人可用'),
    ('当老板不当打字员', '定义目标，让AI执行\n你只做选择，不写过程'),
    ('三选一决策法', '永远让AI给3个方案\n你只负责拍板'),
]
for i, (t, d) in enumerate(cards):
    x = Inches(0.8 + i * 4)
    add_rounded(slide, x, Inches(4.5), Inches(3.5), Inches(2.2), RGBColor(0xf5,0xf5,0xf5))
    add_text(slide, Inches(1.1 + i * 4), Inches(4.7), Inches(3), Inches(0.5),
        t, size=20, color=DARK_TEXT, bold=True)
    add_text(slide, Inches(1.1 + i * 4), Inches(5.3), Inches(3), Inches(1.2),
        d, size=15, color=GRAY)

add_page_num(slide, 2, 9)
add_footer_brand(slide)


# ===================================================================
# SLIDE 3: COZE - 文风系统 (黑底)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)

add_text(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.3),
    '02', size=12, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(0.8), Inches(10), Inches(0.8),
    '文风系统 + 持续迭代', size=36, color=WHITE, bold=True)

add_line(slide, Inches(0.5), Inches(1.8), Inches(1.5), ORANGE)

# 4 cards on dark
cards2 = [
    ('文风提取器', '给AI一篇范文\n→ 自动提取风格\n→ 生成同风格新文'),
    ('对比迭代', 'AI原文 vs 手动修改版\n→ 自动总结规律\n→ 持续优化'),
    ('文章炼金师', '情感金句Skill\n→ 10句精炼输出\n→ 每句都是海报'),
    ('摘要专家', '不是概括全文\n→ 制造信息缺口\n→ 最大化打开率'),
]
for i, (t, d) in enumerate(cards2):
    x = Inches(0.5 + i * 3.15)
    add_rounded(slide, x, Inches(2.5), Inches(2.9), Inches(3.5), CARD_DARK)
    # Orange top accent
    add_line(slide, Inches(0.5 + i * 3.15 + 0.3), Inches(2.8), Inches(1), ORANGE, Pt(2))
    add_text(slide, Inches(0.8 + i * 3.15), Inches(3.1), Inches(2.5), Inches(0.5),
        t, size=20, color=WHITE, bold=True)
    add_text(slide, Inches(0.8 + i * 3.15), Inches(3.8), Inches(2.5), Inches(2),
        d, size=15, color=GRAY)

# Bottom highlight
add_text(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
    '核心：AI不是替你写，是学会你怎么写', size=16, color=ORANGE, bold=True)

add_page_num(slide, 3, 9, dark_bg=True)


# ===================================================================
# SLIDE 4: TRANSITION - 为什么要升级 (白底)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.3),
    '03', size=12, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(0.8), Inches(10), Inches(0.8),
    '为什么要升级到 OpenClaw？', size=36, color=DARK_TEXT, bold=True)

add_line(slide, Inches(0.5), Inches(1.8), Inches(1.5))

# Left: 扣子
add_rounded(slide, Inches(0.5), Inches(2.5), Inches(5.8), Inches(3.5), RGBColor(0xf8,0xf8,0xf8))
add_text(slide, Inches(0.8), Inches(2.7), Inches(5), Inches(0.5),
    '扣子能做的', size=22, color=DARK_TEXT, bold=True)
add_multi(slide, Inches(0.8), Inches(3.4), Inches(5), Inches(2.5),
    [('\u2022  零代码快速搭建', GRAY, False),
     ('\u2022  技能商店生态丰富', GRAY, False),
     ('\u2022  入门门槛极低', GRAY, False),
     ('\u2022  适合标准化流程', GRAY, False)],
    size=17)

# Right: OpenClaw
add_rounded(slide, Inches(6.8), Inches(2.5), Inches(6), Inches(3.5), BLACK)
add_text(slide, Inches(7.1), Inches(2.7), Inches(5.5), Inches(0.5),
    '扣子做不到的', size=22, color=WHITE, bold=True)
add_multi(slide, Inches(7.1), Inches(3.4), Inches(5.5), Inches(2.5),
    [('\u2022  本地执行权限（操作文件/调用API）', GRAY, False),
     ('\u2022  24小时在线（本地部署）', GRAY, False),
     ('\u2022  个性化记忆（跨会话持久）', GRAY, False),
     ('\u2022  知识库联动（Obsidian + 飞书）', GRAY, False)],
    size=17)

# Bottom transition
add_rect(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), ORANGE)
add_text(slide, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.7),
    '扣子让你入门  →  OpenClaw 让你全自动', size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_page_num(slide, 4, 9)
add_footer_brand(slide)


# ===================================================================
# SLIDE 5: OPENCLAW 内容工厂 (黑底)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)

add_text(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.3),
    '04', size=12, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(0.8), Inches(10), Inches(0.8),
    'OpenClaw 内容工厂', size=36, color=WHITE, bold=True)

add_line(slide, Inches(0.5), Inches(1.8), Inches(1.5), ORANGE)

# Flow bar
steps = ['灵感捕获', '选题评分', '大纲生成', '初稿撰写', '爆款标题', 'AI配图', '推送草稿箱', '归档']
for i, step in enumerate(steps):
    x = Inches(0.4 + i * 1.6)
    bg_color = ORANGE if i == 0 else CARD_DARK2
    txt_color = WHITE
    shape = add_rounded(slide, x, Inches(2.5), Inches(1.4), Inches(0.7), bg_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(12)
    p.font.color.rgb = txt_color
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

# Big numbers
nums = [('20', '个触发词'), ('28', '个技能'), ('1', '句话启动')]
for i, (n, label) in enumerate(nums):
    x = Inches(1 + i * 4)
    add_rounded(slide, x, Inches(3.8), Inches(3.2), Inches(2.5), CARD_DARK)
    add_text(slide, x, Inches(4), Inches(3.2), Inches(1.2),
        n, size=56, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.2), Inches(3.2), Inches(0.6),
        label, size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Bottom hint
add_text(slide, Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
    '说「内容工厂发布」四个字 → 初稿自动转HTML → 推送草稿箱 → 归档', size=14, color=ORANGE)

add_page_num(slide, 5, 9, dark_bg=True)


# ===================================================================
# SLIDE 6: 知识库 (白底)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.3),
    '05', size=12, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(0.8), Inches(10), Inches(0.8),
    '知识库：从任何地方获取素材', size=36, color=DARK_TEXT, bold=True)

add_line(slide, Inches(0.5), Inches(1.8), Inches(1.5))

# Three pillars banner
add_rounded(slide, Inches(3.5), Inches(2.2), Inches(6.3), Inches(0.8), BLACK)
add_text(slide, Inches(3.5), Inches(2.25), Inches(6.3), Inches(0.7),
    'Get笔记 + Obsidian + OpenClaw', size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

sources = [
    ('视频 / 音频', '自动转录\nAI摘要 → 同步Obsidian'),
    ('飞书文档', 'API直接读取\n纳入知识库'),
    ('小红书热搜', '浏览器爬取\n结构化数据'),
    ('任意网页', '一句「记一下」\n保存+摘要+同步'),
]
for i, (t, d) in enumerate(sources):
    x = Inches(0.5 + i * 3.15)
    add_rounded(slide, x, Inches(3.5), Inches(2.9), Inches(3), RGBColor(0xf5,0xf5,0xf5))
    add_line(slide, Inches(0.5 + i * 3.15 + 0.3), Inches(3.8), Inches(1), ORANGE, Pt(2))
    add_text(slide, Inches(0.8 + i * 3.15), Inches(4), Inches(2.5), Inches(0.5),
        t, size=20, color=DARK_TEXT, bold=True)
    add_text(slide, Inches(0.8 + i * 3.15), Inches(4.7), Inches(2.5), Inches(1.5),
        d, size=15, color=GRAY)

add_page_num(slide, 6, 9)
add_footer_brand(slide)


# ===================================================================
# SLIDE 7: 数据驱动 + 多平台 (黑底)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)

add_text(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.3),
    '06', size=12, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(0.8), Inches(10), Inches(0.8),
    '数据驱动 + 多平台分发', size=36, color=WHITE, bold=True)

add_line(slide, Inches(0.5), Inches(1.8), Inches(1.5), ORANGE)

# Left: Data
add_rounded(slide, Inches(0.5), Inches(2.5), Inches(6), Inches(4), CARD_DARK)
add_text(slide, Inches(0.8), Inches(2.7), Inches(5), Inches(0.5),
    '数据驱动选题', size=22, color=WHITE, bold=True)

data_items = [
    '小红书热帖爬取 → 批量获取标题、点赞数',
    '内容数据分析 → Excel自动生成图表',
    '选题评分器 → 5维度打分 + 3个标题',
    '灵感捕获 → 链接→评分→金句→选题',
]
for i, item in enumerate(data_items):
    add_text(slide, Inches(0.8), Inches(3.5 + i * 0.7), Inches(5.5), Inches(0.6),
        '\u2022  ' + item, size=15, color=GRAY)

# Right: Multi-platform
add_rounded(slide, Inches(7), Inches(2.5), Inches(5.8), Inches(4), CARD_DARK)
add_text(slide, Inches(7.3), Inches(2.7), Inches(5), Inches(0.5),
    '多平台分发', size=22, color=WHITE, bold=True)

add_multi(slide, Inches(7.3), Inches(3.5), Inches(5), Inches(3),
    [('初稿 (Markdown)', GRAY, False),
     ('    ↓', RGBColor(0x40,0x40,0x40), False),
     ('→ 公众号草稿箱  ✅', ORANGE, True),
     ('→ 小红书（扩展中）', GRAY, False),
     ('→ 知识星球（扩展中）', GRAY, False)],
    size=17, spacing=1.5)

add_page_num(slide, 7, 9, dark_bg=True)


# ===================================================================
# SLIDE 8: 三个心法 (白底)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.3),
    '07', size=12, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(0.8), Inches(10), Inches(0.8),
    '三个心法', size=36, color=DARK_TEXT, bold=True)

add_line(slide, Inches(0.5), Inches(1.8), Inches(1.5))

insights = [
    ('01', '老板思维', '定义目标，让AI执行', '扣子和OpenClaw通用\n你只做决策，不做执行'),
    ('02', '触发词系统', '一个词启动一条流水线', '20个契约词\n「内容工厂发布」= 全自动'),
    ('03', '持续迭代', 'AI越用越懂你', '风格学习 + 教训系统\n37次教训，每次都在进化'),
]
for i, (num, title, subtitle, desc) in enumerate(insights):
    y = Inches(2.3 + i * 1.6)
    # Orange number
    add_text(slide, Inches(0.5), y, Inches(1), Inches(0.5),
        num, size=28, color=ORANGE, bold=True)
    # Divider
    add_line(slide, Inches(1.5), Inches(y + Inches(0.1).inches), Pt(1), RGBColor(0xe0,0xe0,0xe0), Inches(0.6))
    # Title
    add_text(slide, Inches(2), y, Inches(3), Inches(0.5),
        title, size=24, color=DARK_TEXT, bold=True)
    # Subtitle
    add_text(slide, Inches(2), Inches(y + Inches(0.55).inches), Inches(3), Inches(0.4),
        subtitle, size=15, color=GRAY)
    # Description
    add_text(slide, Inches(7), y, Inches(5.5), Inches(1.2),
        desc, size=16, color=GRAY)

add_page_num(slide, 8, 9)
add_footer_brand(slide)


# ===================================================================
# SLIDE 9: Q&A (黑底，和封面呼应)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)

add_text(slide, Inches(0.5), Inches(2), Inches(12), Inches(1.5),
    'Q & A', size=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_line(slide, Inches(5.5), Inches(3.8), Inches(2.3), ORANGE)

add_text(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.8),
    '公众号：赛博番茄', size=24, color=ORANGE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.6),
    '感谢聆听', size=18, color=GRAY, align=PP_ALIGN.CENTER)


# === SAVE ===
out_dir = r'D:\OpenClaw\workspace\output'
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'openclaw-share-v2.pptx')
prs.save(out_path)
print(f'OK: {out_path}')
